# SPDX-FileCopyrightText: Copyright (c) 2023 NVIDIA CORPORATION & AFFILIATES. All rights reserved.
# SPDX-License-Identifier: MIT
#
# Permission is hereby granted, free of charge, to any person obtaining a
# copy of this software and associated documentation files (the "Software"),
# to deal in the Software without restriction, including without limitation
# the rights to use, copy, modify, merge, publish, distribute, sublicense,
# and/or sell copies of the Software, and to permit persons to whom the
# Software is furnished to do so, subject to the following conditions:
#
# The above copyright notice and this permission notice shall be included in
# all copies or substantial portions of the Software.
#
# THE SOFTWARE IS PROVIDED "AS IS", WITHOUT WARRANTY OF ANY KIND, EXPRESS OR
# IMPLIED, INCLUDING BUT NOT LIMITED TO THE WARRANTIES OF MERCHANTABILITY,
# FITNESS FOR A PARTICULAR PURPOSE AND NONINFRINGEMENT. IN NO EVENT SHALL
# THE AUTHORS OR COPYRIGHT HOLDERS BE LIABLE FOR ANY CLAIM, DAMAGES OR OTHER
# LIABILITY, WHETHER IN AN ACTION OF CONTRACT, TORT OR OTHERWISE, ARISING
# FROM, OUT OF OR IN CONNECTION WITH THE SOFTWARE OR THE USE OR OTHER
# DEALINGS IN THE SOFTWARE.
import streamlit as st
import shutil

import time
#import gradio as gr
import argparse
from trt_llama_api import TrtLlmAPI #llama_index does not currently support TRT-LLM. The trt_llama_api.py file defines a llama_index compatible interface for TRT-LLM.
from langchain.embeddings.huggingface import HuggingFaceEmbeddings
from llama_index import LangchainEmbedding, ServiceContext
from llama_index.llms.llama_utils import messages_to_prompt, completion_to_prompt
from llama_index import set_global_service_context
from faiss_vector_storage import FaissEmbeddingStorage

from bs4 import BeautifulSoup
import os
import base64
import requests
from requests.exceptions import HTTPError

#forppt
import json
import re
from io import BytesIO
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

print("I am here")
# Create an argument parser
parser = argparse.ArgumentParser(description='NVIDIA Chatbot Parameters')

# Add arguments
parser.add_argument('--trt_engine_path', type=str, required=False,
                    help="Path to the TensorRT engine.", default=r"llama2\engine")
parser.add_argument('--trt_engine_name', type=str, required=False,
                    help="Name of the TensorRT engine.", default="llama_float16_tp1_rank0.engine")
parser.add_argument('--tokenizer_dir_path', type=str, required=False,
                    help="Directory path for the tokenizer.", default=r"llama2\engine")
parser.add_argument('--embedded_model', type=str, 
                    help="Name or path of the embedded model. Defaults to 'sentence-transformers/all-MiniLM-L6-v2' if "
                         "not provided.",
                    default='sentence-transformers/all-MiniLM-L6-v2')
parser.add_argument('--data_dir', type=str, required=False,
                    help="Directory path for data.", default="./dataset")
parser.add_argument('--verbose', type=bool, required=False,
                    help="Enable verbose logging.", default=False)
# Parse the arguments
args = parser.parse_args()

# Use the provided arguments
trt_engine_path = args.trt_engine_path
trt_engine_name = args.trt_engine_name
tokenizer_dir_path = args.tokenizer_dir_path
embedded_model = args.embedded_model
data_dir = args.data_dir
verbose = args.verbose

# create trt_llm engine object
llm = TrtLlmAPI(
    model_path=trt_engine_path,
    engine_name=trt_engine_name,
    tokenizer_dir=tokenizer_dir_path,
    temperature=0.0,
    max_new_tokens=1024,
    context_window=3900,
    messages_to_prompt=messages_to_prompt,
    completion_to_prompt=completion_to_prompt,
    verbose=False
)

# function for encoding image to base64
def get_image_base64(image_path):
    with open(image_path, "rb") as img_file:
        return base64.b64encode(img_file.read()).decode()

    

# Generating and Publishing new articles to Atlassian Confluence Space
def publish(content, username, confluence_url, confluence_api_token, space_key, use_custom_details=False):
    # if not using custom details, get the default details
    if not use_custom_details:
        username= username
        confluence_url=confluence_url
        confluence_api_token=confluence_api_token
        space_key=space_key

    parent_page_id = None #since we are publising to the space directly
    content = content.lstrip("\n")
    title_end_index = content.find("\n")
    new_page_title = content[:title_end_index].strip()
    new_page_content = content[title_end_index:].strip()

    auth_str = f"{username}:{confluence_api_token}" 

    #set the API endpoint URL
    url = f"{confluence_url}"
    auth_str_encoded = base64.b64encode(auth_str.encode()).decode()
    #set the request headers, inlcuding the API token for authentication
    headers = {
    "Authorization": f"Basic {auth_str_encoded}",
        "Content-Type": "application/json"
    }

    # Set the request payload with the new page information
    data = {
        "type": "page",
        "title": new_page_title,
        "space": {"key": space_key},
        "body": {
            "storage": {
                "value": new_page_content,
                "representation": "storage",
            }
        }
    }
    # If the new page should be a child page, specify the parent page ID
    if parent_page_id:
       data["ancestors"]=[{"type": "page", "id": parent_page_id}]


    try:
        # send post request to create the new page
        response = requests.post(url, headers=headers, json=data)
        # If the response was successful, no Exception will be raised
        response.raise_for_status()
    except HTTPError as http_err:
        # If status code is 400, it might be due to duplicate page title
        if response.status_code == 400:
            st.write(f'HTTP error occurred: {http_err}.')
            st.write(f'Error message: {response.json()["message"].split(".")[-1]}')
        else:
            st.write(f'HTTP error occurred: {http_err}.')
    except Exception as err:
        st.write(f'Other error occurred: {err}.')
    else:
        st.write('Page successfully created.')

#parsing the generated text to json
import re

def parse_format(text):
    slides = []
    slide = {}
    points = []

    lines = text  # Assuming 'text' is a single string, split it into lines.

    for original_line in lines:
        # Preserve the original line with leading spaces
        line = original_line.strip()

        if original_line.startswith("Slide"):
            if slide:
                slide["points"] = points
                slides.append(slide)
                slide = {}
                points = []
            _, subtitle = original_line.split(":", 1)
            slide["subtitle"] = subtitle.strip()
        
        elif original_line.startswith("\t"):
            if original_line.strip().startswith(("*", "-", "•", "+")) or bool(re.match(r"^\d+\.", original_line.strip())):
                points.append(" "*4+original_line.strip()[1:])
        elif original_line.startswith(("*", "-", "•", "+")) or bool(re.match(r"^\d+\.", original_line)):
            points.append(original_line[1:])

    if slide:
        slide["points"] = points
        slides.append(slide)

    return {"slides": slides}


def create_ppt_from_json(json_data):
    """
    Create a PowerPoint presentation from parsed JSON data.
    """
    prs = Presentation()

    # Set slide width and height (16:9 aspect ratio)
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    #Access the slide master
    slide_master = prs.slide_master

    #Define footer properties
    footer_text = "Confidential and Proprietary. © 2024 Karaam Analytics. All rights reserved."
    footer_left = Inches(0.5)
    footer_top = prs.slide_height - Inches(1)
    footer_width = prs.slide_width - Inches(1)
    footer_height = Inches(0.5)

    #get logo
    karaam_logo = "static/logo.jpg"
    conf_logo = "static/conf_logo.jpg"

    logo_width = Inches(1.5)  # Adjust as needed
    logo_x_position = prs.slide_width - logo_width - Inches(0.5)  # Adjust the 0.5 inch offset as needed

    logo_height = Inches(0.5)  # Adjust as needed
    logo_y_position = prs.slide_height - logo_height - Inches(1.0)  # Adjust the 0.5 inch offset as needed


    for slide_data in json_data["slides"]:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        #Adjust the width of the content placeholder
        content.width = Inches(12.5)
        content.top = title.top + title.height + Inches(0.5)  # Move content down by 0.5 inches below the title


        title.text = slide_data["subtitle"]
        # Set the title text color to blue
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 85, 183)  # RGB values for blue

        
        for point in slide_data["points"]:
            p = content.text_frame.add_paragraph()
            p.text = point
            p.level = 0  # Bullet point level
            p.space_after = Pt(14) 
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

        # Add logo to the top right corner
        slide.shapes.add_picture(karaam_logo, logo_x_position, logo_y_position, width=logo_width)

        # Add footer to the slide
        footer_shape = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
        text_frame = footer_shape.text_frame
        p = text_frame.add_paragraph()
        p.text = footer_text
        p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        
    return prs


# # create embeddings model object
embed_model = LangchainEmbedding(HuggingFaceEmbeddings(model_name=embedded_model))
service_context = ServiceContext.from_defaults(llm=llm, embed_model=embed_model) 
set_global_service_context(service_context)

# For content generation
def chatbot(query, query_engine, history):
    response = query_engine.query(query)
    return response

#for chattting with documentation
def chatbot1(query, query_engine, history):
    prompt = f"""
You are an intelligent assistant named KaraamAI. Your job is to provide answers to questions purely based on the context information available to you.
Based on the context available to you and the previous conversation history here in quotes: '{history}', generate a short response to the following question: '{query}'.
Do NOT include the question in the response. Do not say based on the context provided. Just simply provide a short response to the query! 

"""
    response = query_engine.query(prompt)
    return response

def scrape_and_save_articles(email, token, base_url):
    headers = {
        "Authorization": "Basic " + base64.b64encode(f"{email}:{token}".encode()).decode(),
        "Accept": "application/json",
        "Content-Type": "application/json"
    }

    all_pages = []
    limit = 100
    start = 0

    while True:
        response = requests.get(f"{base_url}?limit={limit}&start={start}", headers=headers)
        page_summaries = response.json()

        if not page_summaries["results"]:
            break

        for summary in page_summaries['results']:
            page_id = summary['id']
            response = requests.get(f"{base_url}/{page_id}?expand=body.view", headers=headers)
            page = response.json()

            raw_html = page['body']['view']['value']
            soup = BeautifulSoup(raw_html, 'html.parser')
            text_content = soup.get_text(separator="\n\n")

            # Define the filename and save the content to a text file
            filename = f"article_{page_id}.txt"
            filepath = os.path.join("dataset", filename)
            with open(filepath, "w", encoding="utf-8") as file:
                file.write(f"Title: {page['title']}\n\n{text_content}")

            all_pages.append(filepath)

        start += limit

    return all_pages

def main_streamlit():
    #hide default streamlit footer
    hide_streamlit_style = """
        <style>
            #MainMenu {visibility: hidden;}
            footer {visibility: hidden;}
        </style>

    """

    #give it a ash background
    st.markdown("""
        <style>
            .reportview-container {
                background: #B2BEB5;  
            }
            .main .block-container {
                background: #B2BEB5; 
            }
        </style>
        """,
        unsafe_allow_html=True,
    )

    if 'scraped' not in st.session_state:
        st.session_state.scraped = False

    if not st.session_state['scraped']: 
            logo_conf = get_image_base64("static/conf_logo.jpg")
            st.markdown(f"""
                <div style="text-align: center;">
                    <img src="data:image/jpg;base64,{logo_conf}" alt="Confluence logo" style="height: 100px; padding-left: 5px;">    
                        </div>
            """, unsafe_allow_html=True)
            styled_title_html = """
                <div id="project-title" style="
                    font-size: 2rem;
                    font-weight: bold;
                    text-align: center;
                    margin: 2rem 0;
                    text-shadow: 2px 2px 5px #000000;
                    color: #ffffff;
                    background-color: #1269e4;
                    padding: 0.5rem;
                    border-radius: 0.5rem;
                    overflow: hidden;
                    display: inline-block;
                ">
                    Login...
                </div>
            """

            st.markdown(styled_title_html, unsafe_allow_html=True)
            st.session_state.email_de = st.text_input("Confluence Email", key="email")
            st.session_state.token_de = st.text_input("Confluence API Token", type="password", key="token")
            st.session_state.base_url_de = st.text_input("Confluence Base URL", key="base_url")
            st.session_state.space_key_de = st.text_input("Confluence Space Key", type="password", key="space_key")

            if st.button("submit"):
                with st.spinner('Downloading content...'):
                    if os.path.exists("dataset"):
                        shutil.rmtree("dataset")
                    os.makedirs("dataset", exist_ok=True)
                    try:
                        all_pages = scrape_and_save_articles(st.session_state.email, st.session_state.token, st.session_state.base_url)
                        # load the vectorstore index
                        st.write("Loading the vectorstore index....")
                        faiss_storage = FaissEmbeddingStorage(data_dir=data_dir)
                        st.session_state.query_engine = faiss_storage.get_query_engine()
                        st.success(f"Successfully downloaded and indexed {len(all_pages)} articles.")
                        st.session_state['scraped']=True

                    except Exception as e:
                        st.error(f"Download failed: {e}")
                        st.session_state['scraped']=False

    if 'scraped' in st.session_state and st.session_state['scraped']: 
        print('i proceeded')
        st.markdown(hide_streamlit_style, unsafe_allow_html=True) 
        company_logo = get_image_base64("static/logo.jpg")
        st.markdown(f"""
            <div style="text-align: center;">
                <img src="data:image/jpg;base64,{company_logo}" alt="Karaam Analytics" style="height: 70px; padding-left: 5px;">
                <h4>Chat with Team's Documentation or Contribute</h4>      
                    </div>
        """, unsafe_allow_html=True)


        if 'custom_details' not in st.session_state:
            st.session_state.custom_details = False
            st.session_state.username = ""
            st.session_state.confluence_url = ""
            st.session_state.confluence_api_token = ""
            st.session_state.space_k = ""
            st.session_state.content = "" # Add this line to save the generated content
            st.session_state.content_ppt = "" # Add this line to save the generated ppt content

            if 'input_content' not in st.session_state:
                st.session_state.input_content = "Differences between Karaam threads and Narrative weaves" # Add this line to save the input content
            if 'input_slide' not in st.session_state:
                st.session_state.input_slide = "Differences between Karaam threads and Narrative weaves"
        generate_content = st.radio("What would you like to do?", ("None", "Generate Content", "Generate Powerpoint", "Chat with Documentation"))

        if generate_content == "Generate Content":

            input_content = st.text_area('Enter a title for the article: ', st.session_state.input_content)
            st.session_state.input_content = input_content
            
            input_content = f"""Based on the context provided, write a well formmatted article titled '{input_content}'. You are allowed to be verbose and overly long. Here is an example of the formatting expected. Only Use the structure of this guide ONLY!!: 
            [Title of the article].

            Introduction
            ....
            ....
            ....
            ....
            ....

            [Provide content for the first subheading]
            ....
            ....
            ....
            ....
            ....

            [Provide content for the second subheading]
            ....
            ....
            ....
            ....
            ....

            Conclusion
            ....
            ....
            ....
            ....
            """
            if st.button('Generate'):
                with st.spinner('Generating content...'):
                    response =str(chatbot(input_content, query_engine=st.session_state.query_engine, history="")).replace("</s>", "")
                    print(type(response))
                    if len(response)!=0:
                        st.session_state.content = response
            

            #outside the generate block
            if 'content' in st.session_state and st.session_state['content']:
                st.subheader("Generated content")
                print(type(f"The content type is: {st.session_state.content}"))
                #st.info(st.session_state.content)
                content_container = st.empty()
                #content_container.info(st.session_state.content)

                if st.button('Edit'):
                    st.session_state['edit_mode']=True
                if 'edit_mode' in st.session_state and st.session_state['edit_mode']:
                    edited_content = content_container.text_area("Edit below:", st.session_state.content, key='editable')
                    st.session_state.edited_content = edited_content
                    if st.button('Save changes'):
                        st.session_state['content'] = st.session_state['edited_content']
                        st.session_state['edit_mode'] = False
                        content_container.info(st.session_state.content)
                else:
                    content_container.info(st.session_state.content)
                
                print(st.session_state.content)
                # st.write(st.session_state.content)  
                publish_to_confluence = st.radio("Would you like to publish this content to Confluence?", ("No", "Yes"))
                st.session_state.publish_to_confluence = publish_to_confluence

                if publish_to_confluence == 'Yes':
                    st.sidebar.header("Confluence space details")
                    st.session_state.custom_details = st.sidebar.checkbox('Use custom details', st.session_state.custom_details)
                    if st.session_state.custom_details:
                        st.session_state.username_cus = st.sidebar.text_input("Please enter your confluence email:", st.session_state.username)
                        st.session_state.confluence_url_cus = st.sidebar.text_input("Please enter your confluence space url:", st.session_state.confluence_url)
                        st.session_state.confluence_api_token_cus = st.sidebar.text_input("Please enter your confluence api token:",type='password', value=st.session_state.confluence_api_token)
                        st.session_state.space_key_cus = st.sidebar.text_input("Please enter your confluence space key:", type='password', value=st.session_state.space_k)

                        #are the details complete?
                        if st.session_state.username_cus and st.session_state.confluence_url_cus and st.session_state.confluence_api_token_cus and st.session_state.space_key_cus:
                            if st.button("Publish"):
                                publish(st.session_state.content, use_custom_details=True, username=st.session_state.username_cus, confluence_url=st.session_state.confluence_url_cus, confluence_api_token=st.session_state.confluence_api_token_cus, space_key=st.session_state.space_key_cus)
                    else:
                        st.session_state.username_de = st.session_state.email_de
                        st.session_state.confluence_url_de = st.session_state.base_url_de
                        st.session_state.confluence_api_token_de = st.session_state.token_de
                        st.session_state.space_key_de = st.session_state.space_key_de
                        if st.button("Publish"):
                            publish(st.session_state.content, use_custom_details=False, username=st.session_state.username_de, confluence_url=st.session_state.confluence_url_de, confluence_api_token=st.session_state.confluence_api_token_de, space_key=st.session_state.space_key_de)
                else:
                    # if user selects 'No' do nothing
                    pass
            else:
                pass
        elif generate_content == "Generate Powerpoint":
            #create an input for the query with a default value
            input_slide = st.text_area('Enter presentation title: ', st.session_state.input_slide)
            st.session_state.input_slide = input_slide
            #refine the input by adding additional instructions
            input_slide = f"""Based on the context provided, create a well formatted presentation titled '{input_slide}.' Itemize the points with bullet points. You are allowed to be verbose and overly long. Here is an example of the formatting expected, given the prompt: please create a presentation titled 'Customer Relations is Pivotal to Business Success'.  You can use this as a guide to create your own content. Strictly follow the structure of this guide: 
            Title: Customer Relations is Pivotal to Business Success
            Slide 1: Introduction
            - Customer relations is the key to business success
            - It is important to maintain a good relationship with customers
            
            Slide 2: Customer Relations
            - Customer relations is pivotal to business success
            - Good customer relations leads to customer retention
            - Customer relations is the key to business success
            
            Slide 3: Conclusion
            - Customer relations is pivotal to business success
            - It is important to maintain a good relationship with customers
            """
            print(input_slide)
            if st.button('Generate'):
                #add a spinner
                with st.spinner('Generating presentation...'):
                    response = str(chatbot(input_slide, query_engine=st.session_state.query_engine, history="")).replace("</s>", "")
                    if len(response)!=0:
                        print(response)
                        st.session_state.content_ppt = response
                    else:
                        st.write("Could not generate the ppt based on the available context")
            # Outside the 'Generate' button block
            if 'content_ppt' in st.session_state and st.session_state['content_ppt']:
                st.subheader('Generated Presentation')
                #display the generated content in a white box
                st.info(st.session_state.content_ppt)
                #create radio buttons for downloading the ppt
                download_ppt = st.radio("Do you want to generate and download as ppt?", ("No", "Yes"))
                if download_ppt == "Yes":
                    #Try to parse the generated text in the correct json format for conversion to ppt
                    try:
                        print("I entered")
                        text=st.session_state.content_ppt
                        parsed_data = parse_format(text.split("\n"))
                        print("I did")
                        print(parsed_data)
                        #create a ppt from the parsed json
                        prs = create_ppt_from_json(parsed_data)
                        #convert the ppt to bytes
                        prs_bytes = BytesIO()
                        prs.save(prs_bytes)
                        #save to the session state
                        st.session_state.ppt_data = prs_bytes.getvalue()
                        #create a button to download the ppt
                        st.download_button(
                            label="Download ppt",
                            data = st.session_state.ppt_data,
                            #download the ppt
                            file_name="generated_slides.pptx"
                            )
                    except Exception as e:
                        print(e)
                        st.write("Could not parse the content to json or create the ppt. Try again maybe with a different query.")
                else:
                    # if user selects 'No' do nothing
                    pass
        elif generate_content == "Chat with Documentation":
            small_conf = get_image_base64("static/small_conf.png")
            st.markdown(f"""
                <div style="display: flex; align-items: center;">
                    <img src="data:image/jpg;base64,{small_conf}" alt="Confluence logo" style="height: 40px; margin-left: 10px;">
                    <h2 style="margin: 0; line-height: 50px;">Chat with Documentation</h2>  
                </div>
            """, unsafe_allow_html=True)


            #st.title("Chat with Documentation ")

            if "query_engine" not in st.session_state:
                st.session_state["query_engine"] = faiss_storage.get_query_engine()

            if "messages" not in st.session_state:
                st.session_state["messages"] = []

            if "history" not in st.session_state:
                st.session_state["history"] = ""

            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.markdown(message["content"])

            if prompt := st.chat_input("Ask me anything about your team's documentation."):
                st.session_state.messages.append({"role": "user", "content": prompt})
                with st.chat_message("user",avatar="👴🏻"):
                    st.markdown(prompt)

                with st.chat_message("assistant",avatar="👽"):
                    #if assistant is not a role in session state pass an empty history
                    response = str(chatbot1(prompt, st.session_state.query_engine, st.session_state["history"])).replace("</s>", "")
                    st.write(response)
                st.session_state.messages.append({"role": "assistant", "content": response})
                
                new_history_entry = f'\n"query": {prompt}\n"response": {response}\n'
                st.session_state["history"] += new_history_entry

                #uncomment the line below to keep the history to the last 5 pairs conversations
                #st.session_state["history"] = "\n".join(st.session_state["history"].split("\n")[-10:])  
        else:
            st.write("Choose an option to proceed")

